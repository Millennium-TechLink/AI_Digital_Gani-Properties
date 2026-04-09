from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create a presentation object
prs = Presentation()

# Define Gani Properties Branding
GP_TEAL = RGBColor(15, 58, 61)  # #0F3A3D
GP_GOLD = RGBColor(213, 179, 106)  # #D5B36A
WHITE = RGBColor(255, 255, 255)
DARK_INK = RGBColor(14, 27, 28)

def set_slide_background(slide, color):
    # Set the slide background to a solid color
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle, notes):
    # Slide 0 is the title slide layout
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, GP_TEAL)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Text styling for title slide
    for shape in [title_shape, subtitle_shape]:
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.color.rgb = WHITE
            # Using default serif/sans-serif as fallbacks
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

def add_content_slide(prs, title, bullets, notes):
    # Slide 1 is the title and content layout
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Style title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.color.rgb = GP_TEAL
        paragraph.font.bold = True
    
    # Add bullets
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = ""
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.color.rgb = DARK_INK
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

# --- Slide 1 ---
add_title_slide(prs, 
    "Gani Properties: Digital Transformation & Growth Strategy", 
    "Scaling for the 2026-2027 Bangalore Real Estate Market",
    "Good morning everyone. Today I'm presenting our new digital roadmap. We aren't just launching a website; we are building a growth engine for Gani Properties to dominate the Bangalore market."
)

# --- Slide 2 ---
add_content_slide(prs,
    "The Digital Vision: Why Now?",
    [
        "Mobile-First Market: Over 90% of buyers now start their search on a smartphone.",
        "The Digital Gap: Traditional property sites lack the brand trust and speed we provide.",
        "Our Goal: To own the digital space for residential and agricultural plots in Bangalore.",
        "Outcome: Lower lead acquisition costs and a 24/7 sales presence."
    ],
    "Our customers are online. If we aren't at the top of their search results, we are losing sales to competitors who might have inferior land but better visibility. This strategy fixes that."
)

# --- Slide 3 ---
add_content_slide(prs,
    "Current Capability (Ready Today)",
    [
        "High-Performance UI: Ultra-fast browsing using React & Vite technology.",
        "Interactive property Maps: Real-time location viewing for all plots.",
        "Smart Filtering: Instant search for Residential, Farm, and Agri Lands.",
        "Live Inventory: Showcasing projects in Chinnapannahalli, Doddaballapura, and Sidlaghatta."
    ],
    "We have a foundation that is faster and more interactive than anything our local competitors are using. It's built for speed and ease of use."
)

# --- Slide 4 ---
add_content_slide(prs,
    "SEO Strategy (Visibility)",
    [
        "Target Keywords: 'Residential plots Bangalore', 'BBMP approved plots', 'Farmland in Doddaballapura'.",
        "Local Authority: Dedicated landing pages for North Bangalore (Yelahanka, Devanahalli).",
        "Technical Advantage: Automated metadata and schema markup for Google.",
        "Goal: Reach top 10 organic rankings within the first 3-6 months."
    ],
    "SEO is how we get 'free' customers. By ranking high on Google for terms like 'BBMP approved plots', we ensure that high-intent buyers find Gani Properties first."
)

# --- Slide 5 ---
add_content_slide(prs,
    "Building Trust & Authority",
    [
        "Infrastructure News: Keeping buyers updated on Metro lines and Ring Road developments.",
        "Neighborhood ROI Guides: Providing expert analysis on area growth trends.",
        "Local Trust: Integration with Google Business Profiles and verified location data.",
        "Expertise Positioning: Defining Gani Properties as the 'trusted advisor'."
    ],
    "Trust is the currency of real estate. When we educate the buyer on ROI and infrastructure, they are much more likely to close a deal with us."
)

# --- Slide 6 ---
add_content_slide(prs,
    "Automation & Lead Management",
    [
        "Instant Alerts: Leads sent immediately to Sales Team via WhatsApp/Email.",
        "Smart Lead Filtering: Forms designed to capture high-intent buyers only.",
        "Centralized Admin: Easy dashboard to add/edit/delete properties in seconds.",
        "Zero Leakage: Every inquiry is tracked from start to finish."
    ],
    "We’ve optimized the funnel. A lead can submit an inquiry and have a salesperson calling them back in less than 60 seconds."
)

# --- Slide 7 ---
add_content_slide(prs,
    "Future Scaling (Roadmap 2026-27)",
    [
        "360° Virtual Tours: Enabling NRI/Outstation buyers to visit sites virtually.",
        "AI Recommendation Engine: Automatically matching properties to buyer preference.",
        "Secure Document Portal: Viewing land titles securely online.",
        "Sustainability Focus: Highlighting 'Green' and solar-ready agricultural plots."
    ],
    "This is just the beginning. Our architecture is built to scale into virtual reality tours and AI-driven matching as we grow."
)

# --- Slide 8 ---
add_content_slide(prs,
    "Financial ROI & Cost Efficiency",
    [
        "Fixed Cost Efficiency: ₹0 Monthly Hosting during initial growth phase.",
        "Cost Savings: Digital lead generation is up to 60% cheaper than print ads.",
        "Pay-as-you-Scale: Infrastructure costs only increase as revenue grows.",
        "Market Dominance: Re-investing savings into targeted Google Ads."
    ],
    "We are being financially smart. We've used modern tech to keep our fixed costs at near-zero while we focus our budget on actual marketing and sales."
)

# --- Slide 9 ---
add_content_slide(prs,
    "Ongoing Digital Marketing",
    [
        "Social Media Integration: Automated sharing of new plots to Instagram/Facebook.",
        "Google Maps Presence: Every plot pinned for easy search navigation.",
        "Targeted Newsletters: Keeping our database updated on 'New Arrivals'.",
        "Community Building: Engaging with buyers via area development blogs."
    ],
    "We will have a 24/7 marketing presence that reminds our audience of our value every single week."
)

# --- Slide 10 ---
add_content_slide(prs,
    "Building the Future Together",
    [
        "Immediate Action: Deploy platform to live domain.",
        "Immediate Action: Submit site to Google Search Console.",
        "Immediate Action: Complete full inventory upload via dashboard.",
        "Final Vision: Dominating the digital property landscape by 2026."
    ],
    "Let's launch, scale, and build the future of Gani Properties together. The digital landscape is ready for a leader. Gani Properties is ready to lead."
)

# Save the presentation
output_file = "Gani_Properties_Digital_Strategy.pptx"
prs.save(output_file)

print(f"Presentation successfully saved as {output_file}")
